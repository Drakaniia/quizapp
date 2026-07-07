#!/usr/bin/env python
"""
pptx_to_md.py - Extract every detail from a PPTX lecture into a single
rich .md file. Any AI can read this to generate a comprehensive quiz.

Features:
  - Full text extraction with shape positions and hierarchy
  - Table extraction with cell-level data
  - Image extraction (--images) *and* optional OCR (--ocr)
  - Shape position tracking => layout-aware output
  - Auto-grouping of nearby text with images
  - Slide-level layout description
  - Detected lecture sections
  - AI quiz generation prompt at the end

Usage:
    python script/pptx_to_md.py "files/My Lecture.pptx"
    python script/pptx_to_md.py "files/My Lecture.pptx" --images   # also save images
    python script/pptx_to_md.py "files/My Lecture.pptx" --ocr      # + OCR on images
    python script/pptx_to_md.py "files/My Lecture.pptx" --all      # images + OCR

Output:
    output/<slug>/<slug>.md         # Master markdown — AI-ready
    output/<slug>/images/            # Extracted images (if --images/--all)
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── Optional dependencies ──────────────────────────────────────────────

HAVE_OCR = False
try:
    import pytesseract
    from PIL import Image as PILImage
    import io

    # Quick health check: can tesseract actually be called?
    import subprocess
    subprocess.run(
        ["tesseract", "--version"],
        capture_output=True,
        timeout=5,
    )
    HAVE_OCR = True
except Exception:
    pass

# ── Constants ──────────────────────────────────────────────────────────

EMU_PER_INCH = 914400


def emu_to_in(emu: int) -> float:
    return emu / EMU_PER_INCH


def slugify(name: str) -> str:
    s = name.replace(".pptx", "").replace(".pdf", "").replace(".ppt", "")
    s = re.sub(r"[^a-zA-Z0-9]+", "_", s).strip("_").lower()
    return s

# ── Low-level shape introspection ──────────────────────────────────────


def describe_fill(shape) -> str:
    """Return a colour string for the shape fill."""
    try:
        ft = shape.fill.type
        if ft is None:
            return "none"
        from pptx.enum.dml import MSO_FILL

        if ft == MSO_FILL.SOLID:
            fc = shape.fill.fore_color
            try:
                return "#" + str(fc.rgb)
            except Exception:
                try:
                    return "theme:" + str(fc.theme_color)
                except Exception:
                    return "solid(?)"
        return str(ft)
    except Exception:
        return "unknown"


def extract_runs(paragraph) -> list[dict]:
    """Run-level formatting from a paragraph."""
    runs = []
    for r in paragraph.runs:
        run = {
            "text": r.text,
            "bold": r.font.bold,
            "italic": r.font.italic,
            "size_pt": None,
            "color": None,
            "font": r.font.name,
        }
        try:
            if r.font.size:
                run["size_pt"] = round(r.font.size.pt, 1)
        except Exception:
            pass
        try:
            if r.font.color and r.font.color.rgb:
                run["color"] = "#" + str(r.font.color.rgb)
        except Exception:
            pass
        runs.append(run)
    return runs


def extract_text_frame(tf) -> list[dict]:
    """Paragraph + run structure from a text frame."""
    paras = []
    if tf is None:
        return paras
    for p in tf.paragraphs:
        paras.append(
            {
                "text": p.text,
                "alignment": str(p.alignment) if p.alignment else None,
                "level": p.level,
                "runs": extract_runs(p),
            }
        )
    return paras


def extract_table_data(table) -> dict:
    """Table structure including cell text and spans."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append({
                "text": cell.text.strip(),
                "colspan": int(cell._tc.get("gridSpan", 1)),
                "rowspan": int(cell._tc.get("rowSpan", 1)),
            })
        rows.append(cells)
    return {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "data": rows,
    }


def try_ocr_image(image_blob: bytes) -> str | None:
    """Run Tesseract OCR on an image blob. Returns text or None."""
    if not HAVE_OCR:
        return None
    try:
        img = PILImage.open(io.BytesIO(image_blob))
        text = pytesseract.image_to_string(img, lang="eng").strip()
        return text if text else None
    except Exception:
        return None


# ── Per-shape extraction ───────────────────────────────────────────────


def extract_shape(shape, slide_index: int, shape_idx: int,
                  img_dir: Path | None = None, level: int = 0,
                  do_ocr: bool = False) -> dict:
    """Recursively extract all properties from a single shape.
    Returns a dict with type, text, bounds, fill, children, etc.
    """
    info = {
        "level": level,
        "index": shape_idx,
        "name": shape.name,
        "shape_type": str(shape.shape_type).split(".")[-1] if shape.shape_type else "UNKNOWN",
        "has_text": shape.has_text_frame,
        "text": shape.text_frame.text.strip() if shape.has_text_frame else "",
    }

    # ── Bounds ──
    if hasattr(shape, "left") and shape.left is not None:
        info["bounds"] = {
            "left_emu": int(shape.left),
            "top_emu": int(shape.top),
            "width_emu": int(shape.width),
            "height_emu": int(shape.height),
            "left_in": round(emu_to_in(shape.left), 3),
            "top_in": round(emu_to_in(shape.top), 3),
            "width_in": round(emu_to_in(shape.width), 3),
            "height_in": round(emu_to_in(shape.height), 3),
        }
        # Normalised position relative to slide (0-1)
        # These will be filled in later by build_slide_texts
        info["bounds"]["left_norm"] = None
        info["bounds"]["top_norm"] = None
        info["bounds"]["width_norm"] = None
        info["bounds"]["height_norm"] = None

    # Rotation
    try:
        if shape.rotation:
            info["rotation"] = int(shape.rotation)
    except Exception:
        pass

    # Visuals
    info["fill"] = describe_fill(shape)

    # Text details (run-level)
    if shape.has_text_frame:
        info["text_details"] = extract_text_frame(shape.text_frame)

    # Table
    if shape.has_table:
        info["table"] = extract_table_data(shape.table)

    # Image
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        img_info = {"type": "embedded"}
        try:
            img = shape.image
            ext = img.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            img_info["content_type"] = img.content_type
            img_info["size_bytes"] = len(img.blob)
            img_info["ext"] = ext
            img_info["dpi"] = getattr(shape.image, "dpi", None)

            # Save to disk
            if img_dir:
                fname = "slide{:03d}_img{:03d}.{}".format(slide_index, shape_idx, ext)
                (img_dir / fname).write_bytes(img.blob)
                img_info["file"] = fname

                # OCR
                if do_ocr and HAVE_OCR:
                    ocr_text = try_ocr_image(img.blob)
                    if ocr_text:
                        img_info["ocr_text"] = ocr_text
        except Exception as e:
            img_info["error"] = str(e)

        # Try to get alt-text from XML
        try:
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                  "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
            nvPicPr = shape._element.find(".//p:nvPicPr", ns)
            if nvPicPr is not None:
                cNvPr = nvPicPr.find("p:cNvPr", ns)
                if cNvPr is not None:
                    alt = cNvPr.get("descr", "") or cNvPr.get("title", "")
                    if alt:
                        img_info["alt_text"] = alt
        except Exception:
            pass

        info["image"] = img_info

    # Group children (recursive)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        children = []
        for ci, child in enumerate(shape.shapes):
            children.append(
                extract_shape(child, slide_index, ci, img_dir, level + 1, do_ocr)
            )
        info["children"] = children
        info["child_count"] = len(children)

    # Placeholder type
    if shape.is_placeholder:
        info["placeholder"] = {
            "idx": shape.placeholder_format.idx,
            "type": str(shape.placeholder_format.type),
        }

    # Try to detect shape geometry from XML (for auto shapes)
    if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX):
        try:
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                  "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
            spPr = shape._element.find(".//p:spPr", ns)
            if spPr is None:
                spPr = shape._element.find(".//a:spPr", ns)
            if spPr is not None:
                prstGeom = spPr.find("a:prstGeom", ns)
                if prstGeom is not None:
                    info["preset_geometry"] = prstGeom.get("prst", "rect")
        except Exception:
            pass

    return info


# ── Shape normalisation helpers ────────────────────────────────────────


def normalise_bounds(shape_info: dict, slide_w_emu: int, slide_h_emu: int) -> None:
    """Inject normalised (0-1) positions into a shape info dict (and children)."""
    b = shape_info.get("bounds")
    if b:
        b["left_norm"] = round(b["left_emu"] / slide_w_emu, 4)
        b["top_norm"] = round(b["top_emu"] / slide_h_emu, 4)
        b["width_norm"] = round(b["width_emu"] / slide_w_emu, 4)
        b["height_norm"] = round(b["height_emu"] / slide_h_emu, 4)
    for child in shape_info.get("children", []):
        normalise_bounds(child, slide_w_emu, slide_h_emu)


def describe_region(b: dict) -> str:
    """Return a human-readable region label from normalised bounds."""
    l, t, w, h = b["left_norm"], b["top_norm"], b["width_norm"], b["height_norm"]
    x_zone = "left" if l < 0.33 else "center" if l < 0.66 else "right"
    y_zone = "top" if t < 0.33 else "middle" if t < 0.66 else "bottom"
    size = "small" if w < 0.2 and h < 0.2 else "medium" if w < 0.5 and h < 0.5 else "large"
    return "{}-{} {}".format(y_zone, x_zone, size)


# ── Build structured slide data ────────────────────────────────────────


def build_slide_data(prs: Presentation, img_dir: Path | None,
                     do_ocr: bool) -> list[dict]:
    """Process every slide and return rich dicts with shapes, text, images."""
    slides_data = []
    sw, sh = prs.slide_width, prs.slide_height

    for i, slide in enumerate(prs.slides, 1):
        shapes_list = []
        for si, shape in enumerate(slide.shapes):
            info = extract_shape(shape, i, si, img_dir, 0, do_ocr)
            normalise_bounds(info, sw, sh)
            shapes_list.append(info)

        # Build raw text blocks for AI reading
        text_blocks = _extract_text_blocks(shapes_list)
        raw_text = _build_raw_text(text_blocks)

        slides_data.append({
            "slide_number": i,
            "layout": slide.slide_layout.name,
            "shape_count": len(slide.shapes),
            "shapes": shapes_list,
            "text_blocks": text_blocks,
            "raw_text": raw_text,
        })

    return slides_data


def _extract_text_blocks(shapes: list[dict],
                         level: int = 0) -> list[tuple[int, str, str, dict | None]]:
    """Extract (level, type, text, bounds) from shapes recursively."""
    blocks: list[tuple[int, str, str, dict | None]] = []
    for s in shapes:
        stype = s["shape_type"]
        text = s["text"]
        bounds = s.get("bounds")

        if text:
            blocks.append((level, stype, text, bounds))

        # Table
        if "table" in s:
            tbl = s["table"]
            blocks.append((level, "TABLE",
                           "[Table: {} rows x {} cols]".format(tbl["rows"], tbl["cols"]),
                           bounds))
            for row_data in tbl["data"]:
                row_text = " | ".join(c["text"] for c in row_data)
                blocks.append((level + 1, "ROW", row_text, None))

        # Image
        if "image" in s:
            img = s["image"]
            desc = "[Image]"
            if img.get("file"):
                desc = "[Image: {}]".format(img["file"])
            if img.get("alt_text"):
                desc += " alt=\"{}\"".format(img["alt_text"][:80])
            if img.get("ocr_text"):
                ocr_short = img["ocr_text"][:120].replace("\n", " ")
                desc += " OCR: \"{}\"".format(ocr_short)
            blocks.append((level, "IMAGE", desc, bounds))

            # Show full OCR text in a following block
            if img.get("ocr_text") and len(img["ocr_text"]) > 120:
                blocks.append((level + 1, "OCR", img["ocr_text"], None))

        # Children
        for child in s.get("children", []):
            blocks.extend(_extract_text_blocks([child], level + 1))

    return blocks


def _build_raw_text(blocks: list[tuple[int, str, str, dict | None]]) -> str:
    """Concatenate only readable text (skip structural markers) into one string."""
    lines = []
    for level, stype, text, bounds in blocks:
        if text.startswith("[") and text.endswith("]"):
            continue
        lines.append(text)
    return "\n".join(lines)


# ── Layout analysis ─────────────────────────────────────────────────────


def analyse_slide_layout(shapes: list[dict]) -> list[str]:
    """Describe the spatial arrangement of text and images on a slide."""
    descriptions = []

    # Collect regions
    text_regions = []
    image_regions = []

    def walk(s: dict):
        b = s.get("bounds")
        if not b:
            return
        region = describe_region(b)
        stype = s["shape_type"]
        texts = s.get("text", "")
        if s.get("image"):
            image_regions.append((region, s["image"]))
        if texts:
            # Determine if this looks like a title (large text, top-center)
            is_title = (b["top_norm"] < 0.15 and b["left_norm"] < 0.3
                        and b["width_norm"] > 0.4)
            label = "Title" if is_title else "Text"
            text_regions.append((region, label, texts[:60]))
        for child in s.get("children", []):
            walk(child)

    for s in shapes:
        walk(s)

    if text_regions:
        titles = [t for r, l, t in text_regions if l == "Title"]
        bodies = [t for r, l, t in text_regions if l == "Text"]
        if titles:
            descriptions.append("Title: \"{}\"".format(titles[0]))
        if image_regions:
            descriptions.append("Images: {}".format(len(image_regions)))
        if bodies:
            descriptions.append("Text blocks: {} (e.g. \"{}\")"
                                .format(len(bodies), bodies[0][:50]))
    else:
        descriptions.append("No text content")

    return descriptions


# ── Markdown generation ────────────────────────────────────────────────


def generate_markdown(pptx_path: Path, slides_data: list[dict],
                      prs: Presentation, img_dir: Path | None,
                      do_ocr: bool) -> str:
    """Generate the comprehensive, AI-friendly master markdown."""
    sw, sh = prs.slide_width, prs.slide_height
    file_title = pptx_path.stem
    total_slides = len(slides_data)
    total_shapes = sum(s["shape_count"] for s in slides_data)
    total_images = sum(
        1 for s in slides_data for sh in s["shapes"] if "image" in sh
    )
    total_tables = sum(
        1 for s in slides_data for sh in s["shapes"] if "table" in sh
    )
    total_ocr = sum(
        1 for s in slides_data for sh in s["shapes"]
        if "image" in sh and sh["image"].get("ocr_text")
    )

    lines: list[str] = []

    # ── HEADER ──
    lines.append("# PPTX Lecture Extraction: {}".format(file_title))
    lines.append("")
    lines.append("> **Purpose:** Full-text + layout + image extraction from the lecture.")
    lines.append("> Feed this file to an AI to generate a comprehensive college-level quiz.")
    lines.append("")
    lines.append("## File Overview")
    lines.append("")
    lines.append("| Property | Value |")
    lines.append("|---|---|")
    lines.append("| File | `{}` |".format(pptx_path.name))
    lines.append("| Slides | {} |".format(total_slides))
    lines.append("| Shapes | {} |".format(total_shapes))
    lines.append("| Images | {} |".format(total_images))
    if do_ocr:
        lines.append("| Images OCR-ed | {} |".format(total_ocr))
    lines.append("| Tables | {} |".format(total_tables))
    lines.append("| Dimensions | {:.2f}\" x {:.2f}\" |".format(
        emu_to_in(sw), emu_to_in(sh)))
    if img_dir:
        lines.append("| Image dir | `{}` |".format(img_dir))
    lines.append("")
    lines.append("---")
    lines.append("")

    # ── SLIDE OVERVIEW ──
    lines.append("## Slide Overview")
    lines.append("")
    lines.append("| # | Layout | Shapes | Images | Tables | Key Text |")
    lines.append("|---|---|---|---|---|---|")
    for s in slides_data:
        n_images = sum(1 for sh in s["shapes"] if "image" in sh)
        n_tables = sum(1 for sh in s["shapes"] if "table" in sh)
        preview = s["raw_text"][:80] if s["raw_text"] else "(empty)"
        lines.append("| {} | {} | {} | {} | {} | {} |".format(
            s["slide_number"], s["layout"], s["shape_count"],
            n_images, n_tables, preview.replace("|", "/")))
    lines.append("")

    # ── LECTURE SECTIONS ──
    sections = detect_lecture_sections(slides_data)
    if sections:
        lines.append("## Detected Lecture Sections")
        lines.append("")
        lines.append("| Section | Slides | Topic |")
        lines.append("|---|---|---|")
        for sec in sections:
            lines.append("| {} | {}-{} | {} |".format(
                sec["label"], sec["start"], sec["end"], sec["topic"]))
        lines.append("")
        lines.append("---")
        lines.append("")

    # ── SLIDE-BY-SLIDE ──
    lines.append("## Slide-by-Slide Full Content")
    lines.append("")
    lines.append("---")
    lines.append("")

    for s in slides_data:
        sn = s["slide_number"]

        # Layout description
        layout_desc = analyse_slide_layout(s["shapes"])
        layout_str = "; ".join(layout_desc)
        lines.append("### Slide {} --- `{}`".format(sn, s["layout"]))
        lines.append("")
        lines.append("**Layout:** {}".format(layout_str))
        lines.append("")
        lines.append("**Shapes:** {}".format(s["shape_count"]))
        lines.append("")

        # Tables
        for sh in s["shapes"]:
            if "table" in sh:
                tbl = sh["table"]
                lines.append("#### Table ({}x{})".format(tbl["rows"], tbl["cols"]))
                lines.append("")
                for row_data in tbl["data"]:
                    lines.append("| {} |".format(" | ".join(c["text"] for c in row_data)))
                lines.append("")

        # Images
        images_on_slide = [sh for sh in s["shapes"] if "image" in sh]
        if images_on_slide:
            lines.append("#### Images on this slide")
            lines.append("")
            for img_sh in images_on_slide:
                img = img_sh["image"]
                b = img_sh.get("bounds", {})
                pos = "({:.1f}\", {:.1f}\")".format(
                    b.get("left_in", 0), b.get("top_in", 0))
                size = "{:.1f}\" x {:.1f}\"".format(
                    b.get("width_in", 0), b.get("height_in", 0))
                fname = img.get("file", "?")
                alt = img.get("alt_text", "")
                ocr = img.get("ocr_text", "")

                lines.append("- **`{}`** at {}, {} x {}".format(fname, pos, size, img.get("ext", "?")))
                if img.get("size_bytes"):
                    lines.append("  - Size: {:.1f} KB".format(img["size_bytes"] / 1024))
                if alt:
                    lines.append("  - Alt text: \"{}\"".format(alt))
                if ocr:
                    lines.append("  - OCR text:")
                    for ocr_line in ocr.split("\n"):
                        ocr_line = ocr_line.strip()
                        if ocr_line:
                            lines.append("    > {}".format(ocr_line))
                lines.append("")

        # Full text blocks with position context
        lines.append("#### Extracted Text (with positions)")
        lines.append("")
        lines.append("```")
        for level, stype, text, bounds in s["text_blocks"]:
            prefix = "  " * level
            pos_str = ""
            if bounds:
                pos_str = " [{:.1f}\",{:.1f}\"]".format(
                    bounds["left_in"], bounds["top_in"])
            if text.startswith("[") and text.endswith("]"):
                lines.append("{}{} {}".format(prefix, text, pos_str))
            else:
                lines.append("{}{} {}".format(prefix, text, pos_str))
        lines.append("```")
        lines.append("")

        # Raw continuous text
        if s["raw_text"]:
            lines.append("#### Continuous Text (for AI)")
            lines.append("")
            lines.append(s["raw_text"])
            lines.append("")

        lines.append("---")
        lines.append("")

    # ── QUIZ GENERATION PROMPT ──
    lines.append("## AI Quiz Generation Instructions")
    lines.append("")
    lines.append("Using the lecture content above, generate a college-level quiz:")
    lines.append("")
    lines.append("1. **25 questions** total")
    lines.append("2. **Difficulty mix:** ~10 intermediate, ~10 advanced, ~5 expert")
    lines.append("3. **Question types:**")
    lines.append("   - Factual recall (from explicit slide text)")
    lines.append("   - Application (real-world scenarios)")
    lines.append("   - Analytical (compare/contrast, predict outcomes)")
    lines.append("   - Integrative (combine concepts from multiple slides)")
    lines.append("   - **Diagram-based** (use Image positions and OCR text)")
    lines.append("4. **Each question must have:**")
    lines.append("   - Clear question stem")
    lines.append("   - 4 answer choices (A, B, C, D)")
    lines.append("   - One correct answer")
    lines.append("   - Short explanation (1-3 sentences)")
    lines.append("   - Difficulty label (advanced or expert)")
    lines.append("")
    lines.append("```typescript")
    lines.append("interface Question {")
    lines.append("  id: number;")
    lines.append("  question: string;")
    lines.append("  options: string[];")
    lines.append("  answer: string;")
    lines.append("  explanation: string;")
    lines.append("  difficulty: 'advanced' | 'expert';")
    lines.append("}")
    lines.append("```")
    lines.append("")
    lines.append("Use slide text as primary truth. For diagrams, use the Image positions")
    lines.append("and OCR text (if available) to infer what the diagram shows.")
    lines.append("")

    return "\n".join(lines)


# ── Lecture section detection ──────────────────────────────────────────


def detect_lecture_sections(slides_data: list[dict]) -> list[dict]:
    """Group slides into sections based on topic changes in extracted text."""
    sections: list[dict] = []
    current_label = None
    current_topic = None
    section_start = 1

    for s in slides_data:
        hint = s["raw_text"][:60] if s["raw_text"] else ""
        words = hint.split()
        if 1 <= len(words) <= 8 and hint.isupper() and len(hint) > 3:
            if current_label:
                sections.append({
                    "label": current_label,
                    "start": section_start,
                    "end": s["slide_number"] - 1,
                    "topic": current_topic,
                })
            current_label = "Section {}".format(len(sections) + 1)
            current_topic = hint
            section_start = s["slide_number"]
        elif current_label is None and s["raw_text"]:
            current_label = "Introduction"
            current_topic = s["raw_text"][:50]
            section_start = 1

    if current_label and slides_data:
        sections.append({
            "label": current_label,
            "start": section_start,
            "end": slides_data[-1]["slide_number"],
            "topic": current_topic or "",
        })

    return sections


# ── Main pipeline ──────────────────────────────────────────────────────


def process(pptx_path: Path, output_dir: Path,
            extract_images: bool, do_ocr: bool) -> Path:
    """Run the full extraction pipeline. Returns path to the .md file."""
    print("[PROCESS] " + pptx_path.name)
    if do_ocr and not HAVE_OCR:
        print("[WARN] --ocr requested but Tesseract not found. Install Tesseract")
        print("       (https://github.com/UB-Mannheim/tesseract/wiki) then")
        print("       point pytesseract.pytesseract.tesseract_cmd to tesseract.exe")
        do_ocr = False

    prs = Presentation(str(pptx_path))
    slug = slugify(pptx_path.stem)
    output_dir.mkdir(parents=True, exist_ok=True)

    img_dir = output_dir / "images" if extract_images else None
    if img_dir:
        img_dir.mkdir(parents=True, exist_ok=True)

    slides_data = build_slide_data(prs, img_dir, do_ocr)

    # Write markdown
    md = generate_markdown(pptx_path, slides_data, prs, img_dir, do_ocr)
    md_path = output_dir / "{}.md".format(slug)
    md_path.write_text(md, encoding="utf-8")

    # Write JSON summary
    json_path = output_dir / "{}_data.json".format(slug)
    summary = {
        "file": pptx_path.name,
        "slides": len(slides_data),
        "slide_width_in": round(emu_to_in(prs.slide_width), 2),
        "slide_height_in": round(emu_to_in(prs.slide_height), 2),
        "shape_count": sum(s["shape_count"] for s in slides_data),
        "image_count": sum(
            1 for s in slides_data for sh in s["shapes"] if "image" in sh
        ),
        "table_count": sum(
            1 for s in slides_data for sh in s["shapes"] if "table" in sh
        ),
        "ocr_used": do_ocr and HAVE_OCR,
    }
    json_path.write_text(json.dumps(summary, indent=2), encoding="utf-8")

    # Summary stats
    img_count = summary["image_count"]
    slides_with_text = sum(1 for s in slides_data if s["raw_text"])

    print("  [OK] Master markdown: " + str(md_path))
    print("  [OK] JSON summary: " + str(json_path))
    print("")
    print("[SUMMARY]")
    print("  Slides: {} ({} with text)".format(len(slides_data), slides_with_text))
    print("  Total shapes: {}".format(summary["shape_count"]))
    print("  Tables: {}".format(summary["table_count"]))
    print("  Images: {}".format(img_count))
    if do_ocr:
        print("  OCR: enabled")
    print("  Output: " + str(output_dir.resolve()))

    return md_path


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convert PPTX lectures into AI-ready markdown for quiz generation."
    )
    parser.add_argument("pptx", nargs="?",
                        default="files/7. Biological Macromolecules.pptx",
                        help="Path to the PPTX file")
    parser.add_argument("-o", "--output", default=None,
                        help="Output directory (default: output/<slug>/)")
    parser.add_argument("-i", "--images", action="store_true",
                        help="Extract embedded images to output/images/")
    parser.add_argument("--ocr", action="store_true",
                        help="Run OCR on images (requires Tesseract)")
    parser.add_argument("--all", action="store_true",
                        help="Shortcut for --images --ocr")
    args = parser.parse_args()

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print("[ERROR] {} not found.".format(pptx_path))
        sys.exit(1)

    do_images = args.images or args.all
    do_ocr = args.ocr or args.all
    output_dir = Path(args.output) if args.output else Path("output") / slugify(pptx_path.stem)

    md_path = process(pptx_path, output_dir, do_images, do_ocr)

    print("")
    print("[NEXT STEP] Feed the markdown file to an AI to generate a quiz!")
    print("   " + str(md_path.resolve()))


if __name__ == "__main__":
    main()
